from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Add a title slide (Slide 1)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Industrial Internship Presentation"
subtitle.text = "INTERNSHIP DETAIL\nAMPARO SECURE TECH, BHILWARA\n\nDURATION\nFEBRUARY - MAY 2026\n\nDEEPESH DHANWANI\nCYBER SECURITY ENGINEER & TEAM LEAD\nCOMPUTER SCIENCE & ENGINEERING\n\nSANGAM UNIVERSITY, BHILWARA"

# Slide 2: Index
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Index"
tf = body_shape.text_frame
tf.text = "1. About the Organization"
p = tf.add_paragraph()
p.text = "2. Why I Chose Amparo Secure Tech"
p = tf.add_paragraph()
p.text = "3. Company's Products & Services"
p = tf.add_paragraph()
p.text = "4. My Role: Cyber Security Engineer & Team Lead"
p = tf.add_paragraph()
p.text = "5. Work Processes, Tools & Technologies Used"
p = tf.add_paragraph()
p.text = "6. Key Projects & Contributions"
p = tf.add_paragraph()
p.text = "7. Learning Outcomes & Skills Acquired"
p = tf.add_paragraph()
p.text = "8. Conclusion"

# Define slide data
slides_data = [
    {
        "title": "1. About the Organization",
        "bullets": [
            "Amparo Secure Tech is an innovative technology firm specializing in high-end surveillance and advanced artificial intelligence solutions.",
            "Mission: To enhance public safety and corporate security by integrating modern Deep Learning and Computer Vision models.",
            "Core Focus: Developing scalable, automated IoT systems, facial recognition modules, and real-time threat detection platforms."
        ]
    },
    {
        "title": "2. Why I Chose Amparo Secure Tech",
        "bullets": [
            "Pioneering Technology: Opportunity to work extensively with cutting-edge AI, Deep Learning, and real-time Computer Vision.",
            "Practical Cyber Security Application: Blended traditional networking/security domains with AI safely.",
            "Leadership Opportunity: Dynamic environment to step into a Team Lead role, managing implementations and coordination.",
            "Impact-Driven Work: Chance to build real-world security products that protect human lives and digital assets."
        ]
    },
    {
        "title": "3. Company's Products & Services",
        "bullets": [
            "Smart AI Surveillance Systems: Intelligent dashboards for live multi-camera monitoring.",
            "Automated Biometric Solutions: High-speed Facial Recognition used for access control and automated attendance.",
            "Hazard & Threat Detection: Advanced deep learning modules to detect fires, unauthorized weapons, and safety violations.",
            "Traffic & Vehicle Analytics: Automatic Number Plate Recognition (ANPR) and vehicle classification."
        ]
    },
    {
        "title": "4. My Role: Cyber Security Engineer & Team Lead",
        "bullets": [
            "Secured video feeds (RTSP/IP Cameras) and local data storage.",
            "Ensured privacy and data integrity for biometric databases.",
            "Led the development lifecycle of the Smart AI Monitoring Dashboard.",
            "Coordinated tasks among team members and resolved complex streaming bugs.",
            "Managed project codebase integrating OpenCV, YOLOv8, and Flask."
        ]
    },
    {
        "title": "5. Work Processes, Tools & Technologies Used",
        "bullets": [
            "Programming & Frameworks: Python, Flask, HTML/CSS/JS.",
            "AI / Models: YOLOv8 (Tracking/Detection), YuNet & SFace (Facial Recognition).",
            "Libraries/DB: OpenCV, SQLite, CSV Logs.",
            "Work Process: Agile environment, modular components, version control (Git) for rapid deployment."
        ]
    },
    {
        "title": "6. Key Projects & Contributions",
        "bullets": [
            "Project: Smart AI Monitoring System - Built a multi-modal web dashboard for unified security.",
            "Engineered low-latency background threads for uninterrupted camera streaming.",
            "Integrated YOLOv8 to flag weapons, lack of helmets, and fire hazards in real-time.",
            "Developed an automated check-in/out logging system that exports daily reports.",
            "Created a dynamic Model Control Panel UI to toggle AI modules and optimize server performance."
        ]
    },
    {
        "title": "7. Learning Outcomes & Skills Acquired",
        "bullets": [
            "Mastered deployment of heavy deep-learning models into production without freezing the server.",
            "Gained strong hands-on experience connecting asynchronous backend video processing with a dynamic web frontend.",
            "Learned how to effectively manage real-time threading issues and efficiently lead technical deliverables.",
            "Deepened understanding of secure data flow when handling sensitive biometric and surveillance data locally."
        ]
    },
    {
        "title": "8. Conclusion",
        "bullets": [
            "The internship was a transformative experience bridging theoretical computer science with high-stakes security applications.",
            "Successfully leading the Smart AI Monitoring project resulted in a functional, state-of-the-art surveillance system.",
            "The skills acquired in team leadership, cyber security, and AI provide a highly competitive edge for future professional endeavors."
        ]
    }
]

# Create slides dynamically
for slide_data in slides_data:
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = slide_data["title"]
    
    tf = body_shape.text_frame
    tf.text = slide_data["bullets"][0]
    
    for bullet in slide_data["bullets"][1:]:
        p = tf.add_paragraph()
        p.text = bullet

# Save the PowerPoint file
prs.save('Internship_Presentation_Amparo.pptx')
print("Presentation successfully generated: Internship_Presentation_Amparo.pptx")
